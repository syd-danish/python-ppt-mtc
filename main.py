from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import requests
import os
from urllib.parse import urlparse
import matplotlib.pyplot as plt
from slide_content import VALUATION_DICT,JAWBONE_REVENUE_DATA,SLIDE_CONTENT,IMAGE_URL

#------------------Image URLs-----------------------
downloaded_image = []
for img in IMAGE_URL:
    try:
        filename = os.path.basename(urlparse(img).path)
        if not filename:
            filename = "image.jpg"
        if os.path.exists(filename):
            print(f"{filename} already exists.")
        else:
            response = requests.get(img, timeout=10)
            response.raise_for_status()
            with open(filename, "wb") as f:
                f.write(response.content)
            print(f"Saved {filename} => {img}")
        downloaded_image.append(filename)
    except requests.exceptions.RequestException as e:
        print(f"Failed to download {img}: {e}")

big_jambox=downloaded_image[0]
jambone_up=downloaded_image[1]
jawbone_headset=downloaded_image[2]
jawbone_team=downloaded_image[3]
alexander_asseily=downloaded_image[4]
hosain_rahman=downloaded_image[5]
jawbone_vs_fitbit=downloaded_image[6]
jawbone_lawsuit=downloaded_image[7]
apple_watch=downloaded_image[8]
jawbone_up_series=downloaded_image[9]
jawbone_health_logo=downloaded_image[10]
hist="histogram.png"
thumbnail = "thumbnail.png" #slide 4
#------------------Video URLs-----------------------
JAWBONE_HEADPHONE_AD="https://www.youtube.com/watch?v=DW3TQpz64rA"
#--------------------CONSTANTS----------------------
CONTENT_FONT_SIZE= Pt(16)
#--------------------INITIALIZE---------------------
myPPT=Presentation()
#--------------------FUNCTIONS----------------------
def set_gradient_background(prs):
    for slide in prs.slides:
        fill = slide.background.fill
        fill.gradient()
        stops = fill.gradient_stops
        stops[0].color.rgb = RGBColor(75, 0, 130)
        stops[1].color.rgb = RGBColor(30, 144, 255)
        stops[1].color.brightness = -0.25

def set_global_fonts(prs, title_font, body_font):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if shape.is_placeholder and "title" in shape.placeholder_format.idx.__str__():
                        run.font.name = title_font
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)  # white text
                    else:
                        run.font.name = body_font
                        run.font.color.rgb = RGBColor(230, 230, 230)

def insert_picture_and_caption(slide,image,picture_position,text_position,caption,hl=None):
    pic=slide.shapes.add_picture(image, left=picture_position[0], top=picture_position[1],
                             width=picture_position[2], height=picture_position[3])
    slide_title_box = slide.shapes.add_textbox(left=text_position[0],top=text_position[1],
                                               width=text_position[2],height=text_position[3])
    slide_title_box.text_frame.text = caption
    p = slide_title_box.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].font.size = CONTENT_FONT_SIZE
    if hl is not None:
        pic.click_action.hyperlink.address=hl
        pic.click_action.hyperlink.target_frame = "_blank"
def generate_title_and_content_template(slide_title,slide_text):
    slide=myPPT.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = slide_title
    slide_content = slide.placeholders[1]
    slide_content.text=slide_text
    text_lines=slide_text.splitlines()
    no_of_lines=len(text_lines)
    for line in range(no_of_lines):
        slide_content.text_frame.paragraphs[line].font.size = CONTENT_FONT_SIZE
    return slide
def generate_title_and_subtitle_template(slide_title,slide_subtitle):
    slide = myPPT.slides.add_slide(title_and_subtitle_layout)
    title = slide.placeholders[0]
    title.text = slide_title
    subtitle = slide.placeholders[1]
    subtitle.text = slide_subtitle
    return slide
def generate_section_header_template(slide_title,slide_subtitle):
    slide = myPPT.slides.add_slide(section_header_layout)
    slide.shapes.title.text=slide_title
    slide.placeholders[1].text=slide_subtitle
    return slide
def create_jawbone_chart(data_dict, chart_title, start_year=2010, end_year=2017, width=6,
                         height=3):
    filtered_years = sorted([year for year in data_dict if start_year <= year <= end_year])
    revenue_vals = [data_dict[year]["revenue"] for year in filtered_years]
    funding_vals = [data_dict[year]["funding"] for year in filtered_years]
    fig, ax = plt.subplots(figsize=(width, height))
    ax.plot(filtered_years, revenue_vals, marker='o', color='teal', linewidth=2, label='Revenue ($M)')
    ax.plot(filtered_years, funding_vals, marker='s', color='orange', linewidth=2, label='Funding ($M)')
    ax.set_xlabel("Year")
    ax.set_ylabel("Amount ($M)")
    ax.set_title(chart_title)
    ax.grid(True, alpha=0.3)
    ax.legend()
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    plt.close(fig)
    img_stream.seek(0)
    return img_stream
def create_valuation_chart(start_year, end_year, chart_title,
                                     width=6, height=3,valuation_dict=VALUATION_DICT):
    filtered_years = sorted([year for year in valuation_dict if start_year <= year <= end_year])
    filtered_vals = [valuation_dict[year] for year in filtered_years]
    fig, ax = plt.subplots(figsize=(width, height))
    ax.plot(filtered_years, filtered_vals, marker='o', color='teal', linewidth=2)
    ax.set_xlabel("Year")
    ax.set_ylabel("Valuation ($B)")
    ax.set_title(chart_title)
    ax.grid(True, alpha=0.3)
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    plt.close(fig)
    img_stream.seek(0)
    return img_stream
def create_slide_6():
    slide = myPPT.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.3), Inches(8), Inches(1))
    title_box.text_frame.text = " The Financial Peak of Jawbone"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)

    slide.shapes.add_picture(chart_stream, left=Inches(0), top=Inches(1.2), width=Inches(10), height=Inches(3.8))
    textBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(8), Inches(1.5))
    tf = textBox.text_frame
    tf.word_wrap = True
    for bullet in [
        "• Valuation rises steadily from $0.05B (2004) to $3B peak (2013)",
        "• Partnerships and product launches fuel growth",
        "• Investor backing: Sequoia, Andreessen Horowitz, BlackRock, Kleiner Perkins, Deutsche Telekom, JP Morgan",
        "• Peak Market Share:- Bluetooth speakers: 45%, Fitness tracker: 18%, Bluetooth headphones: 35%"
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.alignment = PP_ALIGN.LEFT
        p.font.size = CONTENT_FONT_SIZE
    return slide
#-----------------------OBTAIN CHARTS-----------------------
chart_stream = create_valuation_chart(
    start_year=2004, end_year=2013,
    chart_title="Jawbone Valuation Growth (2004–2013)") #slide 6
chart_stream_2=create_valuation_chart (start_year=2013,end_year=2017,
chart_title="Jawbone's plummeting fall during the mid 2010s")

r_f_chart=create_jawbone_chart(JAWBONE_REVENUE_DATA,
"Jawbone Revenue and Funding Data") #slide 11

#-------------------------LAYOUTS---------------------------
title_and_subtitle_layout=myPPT.slide_layouts[0]
title_and_content_layout=myPPT.slide_layouts[1]
section_header_layout=myPPT.slide_layouts[2]
two_content_layout=myPPT.slide_layouts[3]
comparison_layout=myPPT.slide_layouts[4]
title_only_layout=myPPT.slide_layouts[5]
blank_layout=myPPT.slide_layouts[6]
content_with_caption_layout=myPPT.slide_layouts[7]
pic_with_caption_layout=myPPT.slide_layouts[8]

def mainloop():
    slides = []
    no_of_slides = 18
    for i in range(1, no_of_slides + 1):
        if i in [6]:
            slides.append(create_slide_6())
        elif i in [1]:
            slides.append(generate_title_and_subtitle_template(SLIDE_CONTENT[i]["title"], SLIDE_CONTENT[i]["subtitle"]))
        elif i in [7]:
            slides.append(generate_section_header_template(SLIDE_CONTENT[i]["title"], SLIDE_CONTENT[i]["subtitle"]))
        elif i in [18]:
            slide18 = myPPT.slides.add_slide(blank_layout)
            title_box = slide18.shapes.add_textbox(left=Inches(1.8), top=Inches(3), width=Inches(10), height=Inches(3))
            title_box.text_frame.text = "Thank You"
            title_box.text_frame.paragraphs[0].font.size = Pt(90)
            slides.append(slide18)
        else:
            slides.append(generate_title_and_content_template(SLIDE_CONTENT[i]["title"], SLIDE_CONTENT[i]["content"]))
    slides[10].shapes.add_picture(r_f_chart, left=Inches(1), top=Inches(4), width=Inches(8), height=Inches(3.3))
    slides[15].shapes.add_picture(chart_stream_2, left=Inches(0.5), top=Inches(3.9), width=Inches(9), height=Inches(3.4))
    slides[14].shapes.add_picture(hist,left=Inches(1),top=Inches(4.35),width=Inches(8),height=Inches(3))
    insert_picture_and_caption(slide=slides[1], image=hosain_rahman,
                               picture_position=[Inches(0.5), Inches(4.5), Inches(4.25), Inches(2.25)],
                               text_position=[Inches(1.4), Inches(6.75), Inches(3), Inches(0.5)],
                               caption="Hosain Rahman, CEO")
    insert_picture_and_caption(slide=slides[1], image=alexander_asseily,
                               picture_position=[Inches(5.35), Inches(4.5), Inches(4.25), Inches(2.25)],
                               text_position=[Inches(5.65), Inches(6.75), Inches(3), Inches(0.5)],
                               caption="Alexander Asseily, Co-Founder")
    insert_picture_and_caption(slide=slides[2], image=jawbone_team,
                               picture_position=[Inches(1.5), Inches(4), Inches(7), Inches(3)],
                               text_position=[Inches(1.5), Inches(7), Inches(3), Inches(0.5)],
                               caption="Hosain Rahman and his colleagues attending the vanity fair")
    insert_picture_and_caption(slide=slides[3], image=thumbnail,
                               picture_position=[Inches(0.5), Inches(4), Inches(4), Inches(2.25)],
                               text_position=[Inches(0.75), Inches(6.25), Inches(3), Inches(0.5)],
                               caption="Video H-Link to advertisement", hl=JAWBONE_HEADPHONE_AD)
    insert_picture_and_caption(slide=slides[3], image=jawbone_headset,
                               picture_position=[Inches(5.35), Inches(4), Inches(4), Inches(2.25)],
                               text_position=[Inches(5.7), Inches(6.25), Inches(3), Inches(0.5)],
                               caption="Jawbone's full bluetooth set")
    insert_picture_and_caption(slide=slides[4], image=big_jambox,
                               picture_position=[Inches(0.5), Inches(4), Inches(4.25), Inches(2.4)],
                               text_position=[Inches(0.5), Inches(6.4), Inches(3), Inches(0.5)],
                               caption="Jambox: Jawbone's Bluetooth Speaker")

    insert_picture_and_caption(slide=slides[4], image=jambone_up,
                               picture_position=[Inches(5.35), Inches(4), Inches(4.25), Inches(2.4)],
                               text_position=[Inches(5.4), Inches(6.4), Inches(3), Inches(0.5)],
                               caption="Jawbone UP: The FitBit Competitor")
    insert_picture_and_caption(slide=slides[7], image=apple_watch,
                               picture_position=[Inches(1.7),Inches(4.5),Inches(1.25),Inches(2)],
                               text_position=[Inches(0.5), Inches(6.5), Inches(3), Inches(0.5)],
                               caption="Apple Watch's health features")
    insert_picture_and_caption(slide=slides[7],image=jawbone_up_series,
                               picture_position=[Inches(4.5),Inches(4.5),Inches(5),Inches(2)],
                               text_position=[Inches(5.5),Inches(6.5),Inches(3),Inches(.5)],
                               caption="Jawbone UP Series: 1-4")
    insert_picture_and_caption(slide=slides[8], image=jawbone_vs_fitbit,
                               picture_position=[Inches(1.5), Inches(4.3), Inches(7), Inches(2.8)],
                               text_position=[Inches(1.55), Inches(7.1), Inches(3), Inches(.5)],
                               caption="FitBit remained as the more reliable and cheaper"
                                " alternative")
    insert_picture_and_caption(slide=slides[12],image=jawbone_health_logo,
                               picture_position=[Inches(5.50), Inches(4.75), Inches(3.5), Inches(2.25)],
                               text_position=[Inches(5.15),Inches(7),Inches(3),Inches(.5)],
                               caption="Jawbone pivots into Jawbone Health")
    insert_picture_and_caption(slide=slides[12],image=jawbone_lawsuit,
                               picture_position=[Inches(0.75), Inches(4.75), Inches(3.5), Inches(2.25)],
                               text_position=[Inches(0.32),Inches(7),Inches(3), Inches(.5)],
                               caption="Jawbone's lawsuit against Fitbit, 2015")

if __name__=="__main__":
    mainloop()
    set_gradient_background(myPPT)
    set_global_fonts(myPPT, title_font="Poppins", body_font="Poppins")
    myPPT.save("MyMTCPresentation.pptx")
